from __future__ import annotations

import json
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from pydantic import BaseModel, ConfigDict, Field, ValidationError


class PresentationFormatError(ValueError):
    pass


class SlideSpec(BaseModel):
    model_config = ConfigDict(extra="ignore", str_strip_whitespace=True)

    title: str = Field(min_length=1, max_length=120)
    subtitle: str = Field(default="", max_length=180)
    key_message: str = Field(default="", max_length=260)
    bullets: list[str] = Field(min_length=1, max_length=8)


class SourceSpec(BaseModel):
    model_config = ConfigDict(extra="ignore", str_strip_whitespace=True)

    title: str = Field(min_length=1, max_length=180)
    url: str = Field(pattern=r"^https?://", max_length=500)


class DeckSpec(BaseModel):
    model_config = ConfigDict(extra="ignore", str_strip_whitespace=True)

    title: str = Field(min_length=1, max_length=140)
    subtitle: str = Field(default="", max_length=220)
    slides: list[SlideSpec] = Field(min_length=1, max_length=20)
    sources: list[SourceSpec] = Field(default_factory=list, max_length=20)


COLORS = {
    "ink": RGBColor(20, 31, 26),
    "dark": RGBColor(19, 62, 50),
    "dark_2": RGBColor(13, 46, 38),
    "accent": RGBColor(42, 160, 120),
    "accent_2": RGBColor(85, 205, 163),
    "pale": RGBColor(232, 245, 239),
    "pale_2": RGBColor(244, 249, 246),
    "border": RGBColor(205, 226, 216),
    "muted": RGBColor(99, 111, 103),
    "white": RGBColor(255, 255, 255),
    "sand": RGBColor(244, 239, 222),
}


def parse_deck_response(content: str) -> DeckSpec:
    cleaned = content.strip()
    fenced = re.fullmatch(r"```(?:json)?\s*(.*?)\s*```", cleaned, re.DOTALL)
    if fenced:
        cleaned = fenced.group(1)
    if not cleaned.startswith("{"):
        start = cleaned.find("{")
        end = cleaned.rfind("}")
        if start >= 0 and end > start:
            cleaned = cleaned[start : end + 1]

    try:
        payload = json.loads(cleaned)
        return DeckSpec.model_validate(payload)
    except (json.JSONDecodeError, ValidationError) as exc:
        raise PresentationFormatError(
            "Copilot returned an invalid PowerPoint slide structure."
        ) from exc


def _set_background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_shape(
    slide,
    shape_type: MSO_SHAPE,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    fill_color: RGBColor,
    line_color: RGBColor | None = None,
):
    shape = slide.shapes.add_shape(
        shape_type,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    return shape


def _add_text_box(
    slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    font_size: int,
    color: RGBColor,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    vertical: MSO_ANCHOR = MSO_ANCHOR.MIDDLE,
):
    shape = slide.shapes.add_textbox(
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = vertical
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def _add_pill(
    slide,
    *,
    left: float,
    top: float,
    text: str,
    fill_color: RGBColor,
    text_color: RGBColor,
) -> None:
    width = max(1.55, min(4.4, 0.11 * len(text) + 0.8))
    pill = _add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left=left,
        top=top,
        width=width,
        height=0.38,
        fill_color=fill_color,
    )
    frame = pill.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = text_color


def _add_footer(slide, slide_number: int, total: int) -> None:
    _add_text_box(
        slide,
        left=0.68,
        top=7.02,
        width=2.4,
        height=0.25,
        text="MY CHAT",
        font_size=9,
        color=COLORS["muted"],
        bold=True,
    )
    _add_text_box(
        slide,
        left=11.8,
        top=7.02,
        width=0.8,
        height=0.25,
        text=f"{slide_number:02d}",
        font_size=9,
        color=COLORS["muted"],
        bold=True,
        align=PP_ALIGN.RIGHT,
    )
    track = _add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        left=3.15,
        top=7.13,
        width=8.2,
        height=0.035,
        fill_color=COLORS["border"],
    )
    track.line.fill.background()
    progress = max(0.25, 8.2 * slide_number / max(total, 1))
    _add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        left=3.15,
        top=7.13,
        width=progress,
        height=0.035,
        fill_color=COLORS["accent"],
    )


def _add_slide_header(
    slide,
    *,
    eyebrow: str,
    title: str,
    subtitle: str,
) -> None:
    _add_pill(
        slide,
        left=0.7,
        top=0.42,
        text=eyebrow,
        fill_color=COLORS["pale"],
        text_color=COLORS["dark"],
    )
    _add_text_box(
        slide,
        left=0.7,
        top=0.9,
        width=11.9,
        height=0.65,
        text=title,
        font_size=25,
        color=COLORS["dark"],
        bold=True,
    )
    if subtitle:
        _add_text_box(
            slide,
            left=0.72,
            top=1.48,
            width=11.65,
            height=0.4,
            text=subtitle,
            font_size=11,
            color=COLORS["muted"],
        )


def _add_key_message(slide, message: str) -> float:
    if not message:
        return 1.95
    banner = _add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left=0.72,
        top=1.95,
        width=11.85,
        height=0.78,
        fill_color=COLORS["dark"],
    )
    frame = banner.text_frame
    frame.clear()
    frame.margin_left = Inches(0.28)
    frame.margin_right = Inches(0.28)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = message
    run.font.name = "Aptos"
    run.font.size = Pt(17)
    run.font.bold = True
    run.font.color.rgb = COLORS["white"]
    return 2.95


def _add_spotlight_layout(slide, bullet: str, top: float) -> None:
    _add_shape(
        slide,
        MSO_SHAPE.OVAL,
        left=0.95,
        top=top + 0.25,
        width=0.72,
        height=0.72,
        fill_color=COLORS["accent"],
    )
    _add_text_box(
        slide,
        left=1.16,
        top=top + 0.38,
        width=0.3,
        height=0.3,
        text="01",
        font_size=11,
        color=COLORS["white"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _add_text_box(
        slide,
        left=1.95,
        top=top,
        width=9.85,
        height=2.6,
        text=bullet[:520],
        font_size=27,
        color=COLORS["ink"],
        bold=True,
    )
    _add_text_box(
        slide,
        left=1.98,
        top=top + 2.65,
        width=8.8,
        height=0.45,
        text="핵심 메시지",
        font_size=11,
        color=COLORS["accent"],
        bold=True,
    )


def _add_card_layout(slide, bullets: list[str], top: float) -> None:
    columns = 2
    rows = 1 if len(bullets) <= 2 else 2
    gap_x = 0.26
    gap_y = 0.22
    card_width = (11.85 - gap_x) / columns
    card_height = min(1.78, (6.55 - top - gap_y * (rows - 1)) / rows)
    for index, bullet in enumerate(bullets[:4]):
        column = index % columns
        row = index // columns
        left = 0.72 + column * (card_width + gap_x)
        card_top = top + row * (card_height + gap_y)
        card = _add_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left=left,
            top=card_top,
            width=card_width,
            height=card_height,
            fill_color=COLORS["pale_2"],
            line_color=COLORS["border"],
        )
        badge = _add_shape(
            slide,
            MSO_SHAPE.OVAL,
            left=left + 0.25,
            top=card_top + 0.25,
            width=0.48,
            height=0.48,
            fill_color=COLORS["accent"],
        )
        badge.text_frame.clear()
        badge.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        paragraph = badge.text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        run = paragraph.add_run()
        run.text = f"{index + 1:02d}"
        run.font.name = "Aptos"
        run.font.size = Pt(9)
        run.font.bold = True
        run.font.color.rgb = COLORS["white"]
        _add_text_box(
            slide,
            left=left + 0.92,
            top=card_top + 0.2,
            width=card_width - 1.18,
            height=card_height - 0.4,
            text=bullet[:300],
            font_size=16 if len(bullet) < 120 else 14,
            color=COLORS["ink"],
            bold=True,
        )
        card.line.width = Pt(1)


def _add_two_column_layout(slide, bullets: list[str], top: float) -> None:
    midpoint = (len(bullets) + 1) // 2
    groups = [bullets[:midpoint], bullets[midpoint:]]
    labels = ["KEY POINTS", "DETAILS"]
    for column, items in enumerate(groups):
        left = 0.72 + column * 6.06
        panel = _add_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left=left,
            top=top,
            width=5.79,
            height=max(2.9, 6.5 - top),
            fill_color=COLORS["pale_2"] if column == 0 else COLORS["sand"],
            line_color=COLORS["border"],
        )
        _add_text_box(
            slide,
            left=left + 0.3,
            top=top + 0.2,
            width=4.9,
            height=0.35,
            text=labels[column],
            font_size=9,
            color=COLORS["accent"],
            bold=True,
        )
        body = slide.shapes.add_textbox(
            Inches(left + 0.32),
            Inches(top + 0.65),
            Inches(5.1),
            Inches(max(2.0, 5.45 - top)),
        )
        frame = body.text_frame
        frame.clear()
        frame.word_wrap = True
        for item_index, item in enumerate(items):
            paragraph = (
                frame.paragraphs[0]
                if item_index == 0
                else frame.add_paragraph()
            )
            paragraph.text = item[:320]
            paragraph.font.name = "Aptos"
            paragraph.font.size = Pt(15)
            paragraph.font.color.rgb = COLORS["ink"]
            paragraph.space_after = Pt(12)
            paragraph.level = 0
        panel.line.width = Pt(1)


def _add_title_slide(presentation: Presentation, deck: DeckSpec) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _set_background(slide, COLORS["dark_2"])
    _add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        left=0,
        top=0,
        width=0.22,
        height=7.5,
        fill_color=COLORS["accent"],
    )
    _add_shape(
        slide,
        MSO_SHAPE.OVAL,
        left=9.7,
        top=-0.65,
        width=4.5,
        height=4.5,
        fill_color=COLORS["dark"],
    )
    _add_shape(
        slide,
        MSO_SHAPE.OVAL,
        left=10.75,
        top=3.55,
        width=2.7,
        height=2.7,
        fill_color=COLORS["accent"],
    )
    _add_pill(
        slide,
        left=0.95,
        top=0.85,
        text="MY CHAT  ·  WEB-RESEARCHED DECK",
        fill_color=COLORS["accent"],
        text_color=COLORS["white"],
    )
    _add_text_box(
        slide,
        left=0.95,
        top=1.7,
        width=8.8,
        height=1.7,
        text=deck.title,
        font_size=35,
        color=COLORS["white"],
        bold=True,
    )
    _add_text_box(
        slide,
        left=0.98,
        top=3.55,
        width=8.7,
        height=0.85,
        text=deck.subtitle or "웹 검색과 첨부 자료를 기반으로 생성한 프레젠테이션",
        font_size=17,
        color=RGBColor(198, 226, 215),
    )
    _add_text_box(
        slide,
        left=0.98,
        top=6.6,
        width=5.4,
        height=0.32,
        text=f"{len(deck.slides):02d} CONTENT SLIDES  ·  SOURCES INCLUDED",
        font_size=10,
        color=RGBColor(158, 205, 187),
        bold=True,
    )


def _add_agenda_slide(
    presentation: Presentation,
    deck: DeckSpec,
    *,
    slide_number: int,
    total: int,
) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _set_background(slide, COLORS["white"])
    _add_slide_header(
        slide,
        eyebrow="OVERVIEW",
        title="Agenda",
        subtitle=deck.title,
    )
    for index, item in enumerate(deck.slides[:8]):
        column = index % 2
        row = index // 2
        left = 0.72 + column * 6.05
        top = 2.05 + row * 1.03
        _add_text_box(
            slide,
            left=left,
            top=top,
            width=0.62,
            height=0.62,
            text=f"{index + 1:02d}",
            font_size=13,
            color=COLORS["accent"],
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        _add_text_box(
            slide,
            left=left + 0.72,
            top=top,
            width=4.9,
            height=0.62,
            text=item.title,
            font_size=15,
            color=COLORS["ink"],
            bold=True,
        )
        _add_shape(
            slide,
            MSO_SHAPE.RECTANGLE,
            left=left + 0.72,
            top=top + 0.67,
            width=4.65,
            height=0.025,
            fill_color=COLORS["border"],
        )
    _add_footer(slide, slide_number, total)


def _add_content_slide(
    presentation: Presentation,
    slide_spec: SlideSpec,
    *,
    content_index: int,
    slide_number: int,
    total: int,
) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _set_background(slide, COLORS["white"])
    _add_slide_header(
        slide,
        eyebrow=f"SECTION {content_index:02d}",
        title=slide_spec.title,
        subtitle=slide_spec.subtitle,
    )
    content_top = _add_key_message(slide, slide_spec.key_message)
    if len(slide_spec.bullets) == 1:
        _add_spotlight_layout(slide, slide_spec.bullets[0], content_top)
    elif len(slide_spec.bullets) <= 4:
        _add_card_layout(slide, slide_spec.bullets, content_top)
    else:
        _add_two_column_layout(slide, slide_spec.bullets, content_top)
    _add_footer(slide, slide_number, total)


def _add_sources_slide(
    presentation: Presentation,
    deck: DeckSpec,
    *,
    slide_number: int,
    total: int,
) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _set_background(slide, COLORS["dark_2"])
    _add_pill(
        slide,
        left=0.72,
        top=0.48,
        text="REFERENCES",
        fill_color=COLORS["accent"],
        text_color=COLORS["white"],
    )
    _add_text_box(
        slide,
        left=0.72,
        top=1.05,
        width=11.7,
        height=0.75,
        text="Sources",
        font_size=28,
        color=COLORS["white"],
        bold=True,
    )
    _add_text_box(
        slide,
        left=0.74,
        top=1.75,
        width=10.4,
        height=0.42,
        text="웹 검색에 사용한 참고 자료입니다.",
        font_size=12,
        color=RGBColor(185, 215, 203),
    )

    sources = deck.sources[:10]
    for index, source in enumerate(sources):
        column = index % 2
        row = index // 2
        left = 0.72 + column * 6.08
        top = 2.35 + row * 0.83
        card = _add_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left=left,
            top=top,
            width=5.78,
            height=0.68,
            fill_color=COLORS["dark"],
            line_color=RGBColor(42, 96, 77),
        )
        frame = card.text_frame
        frame.clear()
        frame.margin_left = Inches(0.2)
        frame.margin_right = Inches(0.2)
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        paragraph = frame.paragraphs[0]
        title_run = paragraph.add_run()
        title_run.text = f"{index + 1:02d}  {source.title[:68]}"
        title_run.font.name = "Aptos"
        title_run.font.size = Pt(10)
        title_run.font.bold = True
        title_run.font.color.rgb = COLORS["white"]
        url_paragraph = frame.add_paragraph()
        url_run = url_paragraph.add_run()
        url_run.text = source.url[:105]
        url_run.hyperlink.address = source.url
        url_run.font.name = "Aptos"
        url_run.font.size = Pt(8)
        url_run.font.color.rgb = COLORS["accent_2"]

    _add_text_box(
        slide,
        left=0.72,
        top=7.02,
        width=2.4,
        height=0.25,
        text="MY CHAT",
        font_size=9,
        color=RGBColor(157, 193, 178),
        bold=True,
    )
    _add_text_box(
        slide,
        left=11.8,
        top=7.02,
        width=0.8,
        height=0.25,
        text=f"{slide_number:02d}",
        font_size=9,
        color=RGBColor(157, 193, 178),
        bold=True,
        align=PP_ALIGN.RIGHT,
    )


def build_presentation(deck: DeckSpec, output_path: Path) -> int:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    has_agenda = len(deck.slides) > 1
    total = 1 + (1 if has_agenda else 0) + len(deck.slides) + (
        1 if deck.sources else 0
    )
    _add_title_slide(presentation, deck)

    slide_number = 2
    if has_agenda:
        _add_agenda_slide(
            presentation,
            deck,
            slide_number=slide_number,
            total=total,
        )
        slide_number += 1

    for content_index, slide_spec in enumerate(deck.slides, start=1):
        _add_content_slide(
            presentation,
            slide_spec,
            content_index=content_index,
            slide_number=slide_number,
            total=total,
        )
        slide_number += 1

    if deck.sources:
        _add_sources_slide(
            presentation,
            deck,
            slide_number=slide_number,
            total=total,
        )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)
    return len(presentation.slides)
