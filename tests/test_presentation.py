from __future__ import annotations

from pptx import Presentation

from my_chat.presentation import DeckSpec, build_presentation


def test_presentation_uses_structured_varied_layouts(tmp_path) -> None:
    deck = DeckSpec.model_validate(
        {
            "title": "My Chat 제품 소개",
            "subtitle": "웹 검색 기반 기능 발표",
            "slides": [
                {
                    "title": "핵심 메시지",
                    "key_message": "모든 답변은 웹 검색으로 근거를 확인합니다.",
                    "bullets": ["신뢰할 수 있는 출처를 함께 제공합니다."],
                },
                {
                    "title": "주요 기능",
                    "subtitle": "사용자 중심 기능 구성",
                    "key_message": "검색, 파일, PPT를 하나의 대화에서 처리합니다.",
                    "bullets": [
                        "의무 웹 검색",
                        "사용자별 파일 첨부",
                        "PPTX 생성과 다운로드",
                        "개인 메모리",
                    ],
                },
                {
                    "title": "운영 원칙",
                    "key_message": "보안과 사용 편의성을 함께 유지합니다.",
                    "bullets": [
                        "4개 고정 사용자",
                        "최초 비밀번호 변경",
                        "파일 소유자 격리",
                        "대화 삭제 시 파일 정리",
                        "라이트/다크 모드",
                        "랩탑 우선 UX",
                    ],
                },
            ],
            "sources": [
                {
                    "title": "GitHub Copilot documentation",
                    "url": "https://docs.github.com/copilot",
                }
            ],
        }
    )
    output = tmp_path / "my-chat.pptx"

    slide_count = build_presentation(deck, output)

    presentation = Presentation(output)
    assert slide_count == 6
    assert len(presentation.slides) == 6
    assert len(presentation.slides[0].shapes) >= 7
    assert len(presentation.slides[1].shapes) >= 10
    assert len(presentation.slides[2].shapes) >= 8
    assert len(presentation.slides[3].shapes) >= 12
    assert len(presentation.slides[4].shapes) >= 6
    assert str(presentation.slides[0].background.fill.fore_color.rgb) == "0D1117"

    card_text_sizes = [
        run.font.size.pt
        for shape in presentation.slides[3].shapes
        if shape.has_text_frame
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
        if "의무 웹 검색" in run.text and run.font.size
    ]
    column_text_sizes = [
        paragraph.font.size.pt
        for shape in presentation.slides[4].shapes
        if shape.has_text_frame
        for paragraph in shape.text_frame.paragraphs
        if "4개 고정 사용자" in paragraph.text and paragraph.font.size
    ]
    assert card_text_sizes == [18]
    assert column_text_sizes == [17]
